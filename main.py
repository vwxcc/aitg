#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from __future__ import annotations

import asyncio
import base64
import html
import json
import logging
import os
import re
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

from aiohttp import ClientSession, ClientTimeout
from aiogram import Bot, Dispatcher, F, Router
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode
from aiogram.exceptions import (
    TelegramBadRequest,
    TelegramForbiddenError,
    TelegramNetworkError,
    TelegramRetryAfter,
)
from aiogram.filters import Command, CommandStart
from aiogram.types import (
    CallbackQuery,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    KeyboardButton,
    Message,
    ReplyKeyboardMarkup,
)
from dotenv import load_dotenv

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


# =============================================================================
# ENV
# =============================================================================

load_dotenv()

BOT_TOKEN = os.getenv("BOT_TOKEN", "").strip()

AI_BASE_URL = os.getenv("AI_BASE_URL", "").strip().rstrip("/")
AI_API_KEY = os.getenv("AI_API_KEY", "").strip()
AI_MODEL_ID = os.getenv("AI_MODEL_ID", "").strip()
AI_MODEL_NAME = os.getenv("AI_MODEL_NAME", "Qwen 3.5 35B").strip()

SYSTEM_PROMPT = os.getenv(
    "SYSTEM_PROMPT",
    "Отвечай полезно, точно и понятно.",
).strip()

MAX_FILE_SIZE_MB = max(
    1,
    int(os.getenv("MAX_FILE_SIZE_MB", "10")),
)

MAX_CONCURRENT_AI_REQUESTS = max(
    1,
    int(os.getenv("MAX_CONCURRENT_AI_REQUESTS", "2")),
)

MAX_CONCURRENT_FILE_PROCESSING = max(
    1,
    int(os.getenv("MAX_CONCURRENT_FILE_PROCESSING", "2")),
)

AI_TIMEOUT_SECONDS = max(
    10,
    int(os.getenv("AI_TIMEOUT_SECONDS", "120")),
)

HISTORY_COMPRESS_WORDS = max(
    1000,
    int(os.getenv("HISTORY_COMPRESS_WORDS", "50000")),
)

MAX_HISTORY_MESSAGES = max(
    10,
    int(os.getenv("MAX_HISTORY_MESSAGES", "100")),
)

DATA_DIR = Path(
    os.getenv("DATA_DIR", "data")
).expanduser()

BROADCAST_DELAY = max(
    0.0,
    float(os.getenv("BROADCAST_DELAY", "0.05")),
)

BROADCAST_CONCURRENCY = max(
    1,
    int(os.getenv("BROADCAST_CONCURRENCY", "5")),
)

LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO").upper()


# =============================================================================
# SUBSCRIPTION / LIMIT ENV
# =============================================================================

def required_positive_int_env(name: str) -> int:
    raw = os.getenv(name, "").strip()

    if not raw:
        raise RuntimeError(
            f"{name} не задан в ENV."
        )

    try:
        value = int(raw)
    except ValueError:
        raise RuntimeError(
            f"{name} должен быть целым числом."
        )

    if value < 0:
        raise RuntimeError(
            f"{name} не может быть отрицательным."
        )

    return value


REFERRAL_BONUS_DAYS = required_positive_int_env(
    "REFERRAL_BONUS_DAYS"
)

FREE_LIMIT = required_positive_int_env(
    "FREE_LIMIT"
)

SUBSCRIPTION_LIMIT = required_positive_int_env(
    "SUBSCRIPTION_LIMIT"
)

LIMIT_PERIOD_HOURS = required_positive_int_env(
    "LIMIT_PERIOD_HOURS"
)

SUPPORT_USERNAME = os.getenv(
    "SUPPORT_USERNAME",
    "",
).strip().lstrip("@")

LIMIT_PERIOD_SECONDS = LIMIT_PERIOD_HOURS * 60 * 60


# =============================================================================
# ADMIN IDS
# =============================================================================

def parse_admin_ids() -> set[int]:
    raw = os.getenv("ADMIN_TELEGRAM_IDS", "")

    result: set[int] = set()

    for value in raw.replace(";", ",").replace("\n", ",").split(","):
        value = value.strip()

        if not value:
            continue

        if value.startswith("+"):
            value = value[1:].strip()

        if value.lstrip("-").isdigit():
            result.add(int(value))

    return result


ADMIN_IDS: set[int] = parse_admin_ids()


# =============================================================================
# LOGGING
# =============================================================================

logging.basicConfig(
    level=getattr(
        logging,
        LOG_LEVEL,
        logging.INFO,
    ),
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
)

log = logging.getLogger("aitg")


# =============================================================================
# CONSTANTS
# =============================================================================

MAX_FILE_SIZE = MAX_FILE_SIZE_MB * 1024 * 1024

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".jpg",
    ".jpeg",
    ".png",
    ".webp",
}

SUPPORTED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "image/jpeg",
    "image/png",
    "image/webp",
}

STOP_CALLBACK = "aitg:stop"

BUSY_TEXT = (
    "⏳ Я уже обрабатываю предыдущий запрос.\n"
    "Подождите немного 🙂"
)

QUEUE_FULL_TEXT = (
    "⏳ Сейчас большая нагрузка.\n"
    "Попробуйте через несколько секунд 🙂"
)

STOPPED_TEXT = "⏹ Запрос остановлен."

SUBSCRIPTION_BUTTON = "👥 Подписка за друга"

SUBSCRIPTION_COMMAND = "subscription"


# =============================================================================
# DIRECTORIES
# =============================================================================

USERS_FILE = DATA_DIR / "users.json"
SUBSCRIPTIONS_FILE = DATA_DIR / "subscriptions.json"
HISTORY_DIR = DATA_DIR / "history"
FILES_DIR = DATA_DIR / "files"

DATA_DIR.mkdir(parents=True, exist_ok=True)
HISTORY_DIR.mkdir(parents=True, exist_ok=True)
FILES_DIR.mkdir(parents=True, exist_ok=True)


# =============================================================================
# HELPERS
# =============================================================================

def safe_json_load(
    path: Path,
    default: Any,
) -> Any:
    if not path.exists():
        return default

    try:
        with path.open(
            "r",
            encoding="utf-8",
        ) as file:
            return json.load(file)

    except Exception:
        log.exception(
            "Failed to read JSON: %s",
            path,
        )
        return default


def safe_json_save(
    path: Path,
    data: Any,
) -> None:
    path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    temp_path = path.with_suffix(
        path.suffix
        + f".{uuid.uuid4().hex}.tmp"
    )

    try:
        with temp_path.open(
            "w",
            encoding="utf-8",
        ) as file:
            json.dump(
                data,
                file,
                ensure_ascii=False,
                indent=2,
            )

        temp_path.replace(path)

    finally:
        try:
            if temp_path.exists():
                temp_path.unlink()
        except Exception:
            pass


def count_words(
    text: str,
) -> int:
    return len(
        re.findall(
            r"\S+",
            text or "",
        )
    )


def history_words(
    history: list[dict[str, Any]],
) -> int:
    total = 0

    for item in history:
        content = item.get(
            "content",
            "",
        )

        if isinstance(content, str):
            total += count_words(content)

        elif isinstance(content, list):
            for part in content:
                if isinstance(part, dict):
                    text = part.get("text")

                    if isinstance(text, str):
                        total += count_words(text)

    return total


# =============================================================================
# TELEGRAM HTML / MARKDOWN / MATH NORMALIZER
# =============================================================================

def _protect_matches(
    text: str,
    pattern: str,
    storage: list[str],
    flags: int = 0,
) -> str:
    regex = re.compile(
        pattern,
        flags,
    )

    def replace(match: re.Match) -> str:
        index = len(storage)

        storage.append(
            match.group(0)
        )

        return f"\x00PROTECTED_{index}\x00"

    return regex.sub(
        replace,
        text,
    )


def _find_matching_brace(
    text: str,
    start: int,
) -> int:
    if start >= len(text) or text[start] != "{":
        return -1

    depth = 0

    for index in range(
        start,
        len(text),
    ):
        char = text[index]

        if char == "{":
            depth += 1

        elif char == "}":
            depth -= 1

            if depth == 0:
                return index

    return -1


def _latex_extract_braced(
    text: str,
    start: int,
) -> tuple[str, int] | None:
    end = _find_matching_brace(
        text,
        start,
    )

    if end < 0:
        return None

    return (
        text[start + 1:end],
        end + 1,
    )


def _latex_clean_text(
    text: str,
) -> str:
    text = re.sub(
        r"\\text\s*\{([^{}]*)\}",
        r"\1",
        text,
    )

    text = re.sub(
        r"\\mathrm\s*\{([^{}]*)\}",
        r"\1",
        text,
    )

    text = re.sub(
        r"\\mathbf\s*\{([^{}]*)\}",
        r"\1",
        text,
    )

    text = re.sub(
        r"\\mathit\s*\{([^{}]*)\}",
        r"\1",
        text,
    )

    text = re.sub(
        r"\\operatorname\s*\{([^{}]*)\}",
        r"\1",
        text,
    )

    return text


LATEX_SYMBOLS = {
    r"\alpha": "α",
    r"\beta": "β",
    r"\gamma": "γ",
    r"\delta": "δ",
    r"\epsilon": "ε",
    r"\varepsilon": "ε",
    r"\theta": "θ",
    r"\lambda": "λ",
    r"\mu": "μ",
    r"\pi": "π",
    r"\rho": "ρ",
    r"\sigma": "σ",
    r"\phi": "φ",
    r"\varphi": "φ",
    r"\omega": "ω",
    r"\Delta": "Δ",
    r"\Sigma": "Σ",
    r"\Omega": "Ω",
    r"\Gamma": "Γ",
    r"\Theta": "Θ",
    r"\Lambda": "Λ",
    r"\Phi": "Φ",
    r"\Psi": "Ψ",
    r"\infty": "∞",
    r"\cdot": "·",
    r"\times": "×",
    r"\div": "÷",
    r"\pm": "±",
    r"\mp": "∓",
    r"\neq": "≠",
    r"\ne": "≠",
    r"\leq": "≤",
    r"\le": "≤",
    r"\geq": "≥",
    r"\ge": "≥",
    r"\approx": "≈",
    r"\equiv": "≡",
    r"\sum": "Σ",
    r"\prod": "Π",
    r"\int": "∫",
    r"\partial": "∂",
    r"\rightarrow": "→",
    r"\to": "→",
    r"\leftarrow": "←",
    r"\Rightarrow": "⇒",
    r"\Leftrightarrow": "⇔",
    r"\in": "∈",
    r"\notin": "∉",
    r"\subset": "⊂",
    r"\subseteq": "⊆",
    r"\cup": "∪",
    r"\cap": "∩",
}


def _latex_replace_fractions(
    text: str,
) -> str:
    pattern = re.compile(
        r"\\(?:frac|dfrac|tfrac)\s*\{"
    )

    while True:
        match = pattern.search(text)

        if not match:
            break

        brace_start = (
            match.end() - 1
        )

        numerator_result = (
            _latex_extract_braced(
                text,
                brace_start,
            )
        )

        if numerator_result is None:
            break

        numerator, after_numerator = (
            numerator_result
        )

        search_from = after_numerator

        while (
            search_from < len(text)
            and text[search_from].isspace()
        ):
            search_from += 1

        if (
            search_from >= len(text)
            or text[search_from] != "{"
        ):
            break

        denominator_result = (
            _latex_extract_braced(
                text,
                search_from,
            )
        )

        if denominator_result is None:
            break

        denominator, after_denominator = (
            denominator_result
        )

        numerator = _latex_to_plain_math(
            numerator
        )
        denominator = _latex_to_plain_math(
            denominator
        )

        replacement = (
            f"({numerator})"
            f"/"
            f"({denominator})"
        )

        text = (
            text[:match.start()]
            + replacement
            + text[after_denominator:]
        )

    return text


def _latex_replace_roots(
    text: str,
) -> str:
    pattern = re.compile(
        r"\\sqrt"
        r"(?:\s*\[([^\]]+)\])?"
        r"\s*\{"
    )

    while True:
        match = pattern.search(text)

        if not match:
            break

        index_text = match.group(1)

        brace_start = (
            match.end() - 1
        )

        content_result = (
            _latex_extract_braced(
                text,
                brace_start,
            )
        )

        if content_result is None:
            break

        content, after = content_result

        content = _latex_to_plain_math(
            content
        )

        if index_text:
            index_text = _latex_to_plain_math(
                index_text
            )

            replacement = (
                f"√[{index_text}]({content})"
            )
        else:
            replacement = (
                f"√({content})"
            )

        text = (
            text[:match.start()]
            + replacement
            + text[after:]
        )

    return text


def _latex_replace_text_commands(
    text: str,
) -> str:
    commands = (
        "text",
        "mathrm",
        "mathbf",
        "mathit",
        "operatorname",
    )

    for command in commands:
        pattern = re.compile(
            rf"\\{command}\s*\{{"
        )

        while True:
            match = pattern.search(text)

            if not match:
                break

            content_result = (
                _latex_extract_braced(
                    text,
                    match.end() - 1,
                )
            )

            if content_result is None:
                break

            content, after = content_result

            replacement = (
                _latex_to_plain_math(
                    content
                )
            )

            text = (
                text[:match.start()]
                + replacement
                + text[after:]
            )

    return text


def _latex_replace_scripts(
    text: str,
) -> str:
    # Сначала сложные скрипты.
    def replace_braced_script(
        match: re.Match,
    ) -> str:
        base = match.group(1)
        operator = match.group(2)
        value = match.group(3)

        value = _latex_to_plain_math(
            value
        )

        if operator == "^":
            return f"{base}^({value})"

        return f"{base}_({value})"

    text = re.sub(
        r"([A-Za-zΑ-Ωα-ω0-9)\]])"
        r"\s*([\^_])"
        r"\s*\{([^{}]*)\}",
        replace_braced_script,
        text,
    )

    # Затем одиночные степени/индексы.
    text = re.sub(
        r"([A-Za-zΑ-Ωα-ω0-9)\]])"
        r"\s*\^\s*([A-Za-z0-9+\-]+)",
        r"\1^(\2)",
        text,
    )

    text = re.sub(
        r"([A-Za-zΑ-Ωα-ω0-9)\]])"
        r"\s*_\s*([A-Za-z0-9+\-]+)",
        r"\1_(\2)",
        text,
    )

    return text


def _latex_replace_accents(
    text: str,
) -> str:
    replacements = {
        r"\bar": "¯",
        r"\overline": "¯",
        r"\hat": "̂",
        r"\vec": "⃗",
        r"\tilde": "̃",
    }

    for command, symbol in replacements.items():
        pattern = re.compile(
            re.escape(command)
            + r"\s*\{([^{}]*)\}"
        )

        text = pattern.sub(
            lambda m: (
                f"{m.group(1)}{symbol}"
            ),
            text,
        )

    return text


def _latex_replace_delimiters(
    text: str,
) -> str:
    text = text.replace(
        r"\left",
        "",
    )

    text = text.replace(
        r"\right",
        "",
    )

    for command in (
        r"\bigg",
        r"\Bigg",
        r"\big",
        r"\Big",
    ):
        text = text.replace(
            command,
            "",
        )

    for command in (
        r"\,",
        r"\:",
        r"\;",
        r"\!",
        r"\quad",
        r"\qquad",
        r"\ ",
    ):
        text = text.replace(
            command,
            " ",
        )

    return text


def _latex_replace_symbols(
    text: str,
) -> str:
    for command in sorted(
        LATEX_SYMBOLS,
        key=len,
        reverse=True,
    ):
        text = text.replace(
            command,
            LATEX_SYMBOLS[command],
        )

    return text


def _latex_replace_sum_like(
    text: str,
) -> str:
    pattern = re.compile(
        r"(Σ|Π|∫)"
        r"(?:_\(([^)]*)\))?"
        r"(?:\^\(([^)]*)\))?"
    )

    def replace(
        match: re.Match,
    ) -> str:
        symbol = match.group(1)
        lower = match.group(2)
        upper = match.group(3)

        if lower and upper:
            return (
                f"{symbol}"
                f"({lower}…{upper})"
            )

        if lower:
            return (
                f"{symbol}"
                f"({lower}…)"
            )

        if upper:
            return (
                f"{symbol}"
                f"(…{upper})"
            )

        return symbol

    return pattern.sub(
        replace,
        text,
    )


def _latex_to_plain_math(
    text: str,
) -> str:
    text = _latex_replace_text_commands(
        text
    )

    text = _latex_replace_fractions(
        text
    )

    text = _latex_replace_roots(
        text
    )

    text = _latex_replace_accents(
        text
    )

    text = _latex_replace_delimiters(
        text
    )

    text = _latex_replace_symbols(
        text
    )

    text = _latex_replace_scripts(
        text
    )

    text = _latex_replace_sum_like(
        text
    )

    # Оставшиеся LaTeX-команды неизвестного назначения
    # не удаляем полностью. Убираем только обратный слэш,
    # чтобы пользователь не получал мусор вида \foo.
    text = re.sub(
        r"\\([A-Za-z]+)",
        r"\1",
        text,
    )

    # Внутри математического выражения фигурные скобки
    # обычно служат группировкой.
    text = text.replace(
        "{",
        "(",
    ).replace(
        "}",
        ")",
    )

    # Слишком много пробелов.
    text = re.sub(
        r"[ \t]+",
        " ",
        text,
    )

    return text.strip()


def _normalize_math_blocks(
    text: str,
) -> str:
    # Блоки $$...$$
    text = re.sub(
        r"\$\$(.*?)\$\$",
        lambda m: (
            _latex_to_plain_math(
                m.group(1)
            )
        ),
        text,
        flags=re.DOTALL,
    )

    # \[...\]
    text = re.sub(
        r"\\\[(.*?)\\\]",
        lambda m: (
            _latex_to_plain_math(
                m.group(1)
            )
        ),
        text,
        flags=re.DOTALL,
    )

    # \(...\)
    text = re.sub(
        r"\\\((.*?)\\\)",
        lambda m: (
            _latex_to_plain_math(
                m.group(1)
            )
        ),
        text,
        flags=re.DOTALL,
    )

    # Одиночные $...$
    text = re.sub(
        r"(?<!\$)"
        r"\$"
        r"(?!\$)"
        r"(.+?)"
        r"(?<!\$)"
        r"\$"
        r"(?!\$)",
        lambda m: (
            _latex_to_plain_math(
                m.group(1)
            )
        ),
        text,
        flags=re.DOTALL,
    )

    return text


def _convert_markdown_to_html(
    text: str,
) -> str:
    # Защищаем уже существующие HTML-теги Telegram,
    # чтобы последующее экранирование не уничтожило их.
    html_tags: list[str] = []

    text = _protect_matches(
        text,
        r"</?(?:b|strong|i|em|u|s|strike|del|code|pre|blockquote)(?:\s[^>]*)?>",
        html_tags,
        flags=re.IGNORECASE,
    )

    # Экранируем пользовательские HTML-символы.
    text = html.escape(
        text,
        quote=False,
    )

    # Жирный.
    text = re.sub(
        r"\*\*(.+?)\*\*",
        r"<b>\1</b>",
        text,
        flags=re.DOTALL,
    )

    text = re.sub(
        r"__(.+?)__",
        r"<b>\1</b>",
        text,
        flags=re.DOTALL,
    )

    # Курсив — только когда * не окружает пробел.
    text = re.sub(
        r"(?<!\*)\*(?!\s)(.+?)(?<!\s)\*(?!\*)",
        r"<i>\1</i>",
        text,
        flags=re.DOTALL,
    )

    text = re.sub(
        r"(?<!_)_(?!\s)(.+?)(?<!\s)_(?!_)",
        r"<i>\1</i>",
        text,
        flags=re.DOTALL,
    )

    # Восстанавливаем ранее существовавшие HTML-теги.
    def restore(
        match: re.Match,
    ) -> str:
        index = int(
            match.group(1)
        )

        if 0 <= index < len(html_tags):
            return html_tags[index]

        return match.group(0)

    text = re.sub(
        r"\x00PROTECTED_(\d+)\x00",
        restore,
        text,
    )

    return text


def normalize_ai_answer(
    text: str,
) -> str:
    """
    Преобразует ответ модели только для отображения в Telegram.

    Исходный AI-текст не изменяется и должен сохраняться в историю.
    """

    if not isinstance(text, str):
        text = str(text)

    if not text:
        return ""

    protected: list[str] = []

    # -------------------------------------------------------------------------
    # 1. Защита fenced code.
    # -------------------------------------------------------------------------

    text = _protect_matches(
        text,
        r"```(?:[^\n`]*)\n.*?```",
        protected,
        flags=re.DOTALL,
    )

    # -------------------------------------------------------------------------
    # 2. Защита inline code.
    # -------------------------------------------------------------------------

    text = _protect_matches(
        text,
        r"`[^`\n]+`",
        protected,
    )

    # -------------------------------------------------------------------------
    # 3. Защита URL.
    #
    # В URL могут быть _, &, ?, =, {}, скобки и т.д.
    # -------------------------------------------------------------------------

    text = _protect_matches(
        text,
        r"https?://[^\s<>\"]+",
        protected,
        flags=re.IGNORECASE,
    )

    # -------------------------------------------------------------------------
    # 4. Математика.
    # -------------------------------------------------------------------------

    text = _normalize_math_blocks(
        text
    )

    # -------------------------------------------------------------------------
    # 5. Markdown -> Telegram HTML.
    # -------------------------------------------------------------------------

    text = _convert_markdown_to_html(
        text
    )

    # -------------------------------------------------------------------------
    # 6. Восстанавливаем защищённые фрагменты.
    #
    # Для code и URL содержимое нужно HTML-экранировать.
    # -------------------------------------------------------------------------

    def restore_protected(
        match: re.Match,
    ) -> str:
        index = int(
            match.group(1)
        )

        if not (
            0 <= index < len(protected)
        ):
            return match.group(0)

        value = protected[index]

        # Fenced code.
        if value.startswith("```"):
            inner = value[3:]

            if inner.endswith("```"):
                inner = inner[:-3]

            # Убираем название языка из первой строки.
            first_newline = inner.find("\n")

            if first_newline >= 0:
                language = inner[:first_newline].strip()

                if language:
                    inner = inner[
                        first_newline + 1:
                    ]

            return (
                "<pre><code>"
                + html.escape(
                    inner,
                    quote=False,
                )
                + "</code></pre>"
            )

        # Inline code.
        if value.startswith("`") and value.endswith("`"):
            inner = value[1:-1]

            return (
                "<code>"
                + html.escape(
                    inner,
                    quote=False,
                )
                + "</code>"
            )

        # URL.
        if re.match(
            r"https?://",
            value,
            flags=re.IGNORECASE,
        ):
            return html.escape(
                value,
                quote=False,
            )

        return html.escape(
            value,
            quote=False,
        )

    text = re.sub(
        r"\x00PROTECTED_(\d+)\x00",
        restore_protected,
        text,
    )

    # -------------------------------------------------------------------------
    # 7. Telegram HTML не принимает произвольные HTML-теги.
    #
    # Наша конвертация создаёт только разрешённые теги.
    # Но если модель вернула неизвестные теги, они уже были экранированы.
    # -------------------------------------------------------------------------

    return text.strip()


def split_long_text(
    text: str,
    max_length: int = 3900,
) -> list[str]:
    """
    Делит текст на Telegram-сообщения.

    Старается не разрывать HTML-теги.
    Для больших pre/code-блоков сохраняет баланс тегов.
    """

    if len(text) <= max_length:
        return [text]

    parts: list[str] = []
    remaining = text

    while len(remaining) > max_length:
        cut = remaining.rfind(
            "\n",
            0,
            max_length,
        )

        if cut < max_length // 2:
            cut = remaining.rfind(
                " ",
                0,
                max_length,
            )

        if cut < max_length // 2:
            cut = max_length

        # Не разрываем HTML-тег.
        tag_start = remaining.rfind(
            "<",
            0,
            cut,
        )

        tag_end = remaining.rfind(
            ">",
            0,
            cut,
        )

        if tag_start > tag_end:
            cut = tag_start

        if cut <= 0:
            cut = max_length

        candidate = remaining[:cut].strip()

        if not candidate:
            cut = max_length
            candidate = remaining[:cut].strip()

        parts.append(candidate)

        remaining = remaining[cut:].strip()

    if remaining:
        parts.append(remaining)

    return parts


def image_to_data_url(
    path: Path,
    mime_type: str,
) -> str:
    encoded = base64.b64encode(
        path.read_bytes()
    ).decode("ascii")

    return (
        f"data:{mime_type};base64,{encoded}"
    )


def get_file_extension(
    file_name: str | None,
) -> str:
    if not file_name:
        return ""

    return Path(
        file_name
    ).suffix.lower()


def utc_now_timestamp() -> float:
    return time.time()


def format_datetime(
    timestamp: float,
) -> str:
    try:
        dt = datetime.fromtimestamp(
            timestamp,
            tz=timezone.utc,
        ).astimezone()

        return dt.strftime(
            "%d.%m.%Y %H:%M"
        )

    except Exception:
        return "неизвестно"


def format_remaining_subscription(
    expiry: float,
) -> str:
    remaining = (
        expiry
        - utc_now_timestamp()
    )

    if remaining <= 0:
        return "Подписка неактивна"

    total_minutes = max(
        1,
        int(remaining / 60),
    )

    days = total_minutes // (
        24 * 60
    )

    hours = (
        total_minutes
        % (24 * 60)
    ) // 60

    minutes = (
        total_minutes
        % 60
    )

    parts: list[str] = []

    if days:
        parts.append(
            f"{days} д."
        )

    if hours:
        parts.append(
            f"{hours} ч."
        )

    if minutes and len(parts) < 2:
        parts.append(
            f"{minutes} мин."
        )

    return " ".join(parts)


# =============================================================================
# JSON STORAGE
# =============================================================================

class JSONStore:
    def __init__(self):
        self.users_file = USERS_FILE
        self.subscriptions_file = (
            SUBSCRIPTIONS_FILE
        )
        self.history_dir = HISTORY_DIR

        self.users_lock = asyncio.Lock()
        self.subscription_lock = (
            asyncio.Lock()
        )

    # -------------------------------------------------------------------------
    # USERS
    # -------------------------------------------------------------------------

    async def get_users(
        self,
    ) -> list[int]:

        async with self.users_lock:
            data = await asyncio.to_thread(
                safe_json_load,
                self.users_file,
                [],
            )

        result: list[int] = []

        if isinstance(data, list):
            for value in data:
                try:
                    result.append(int(value))
                except (
                    TypeError,
                    ValueError,
                ):
                    continue

        elif isinstance(data, dict):
            raw_users = data.get(
                "users",
                [],
            )

            if isinstance(
                raw_users,
                list,
            ):
                for value in raw_users:
                    try:
                        result.append(
                            int(value)
                        )
                    except (
                        TypeError,
                        ValueError,
                    ):
                        continue

        return list(
            dict.fromkeys(result)
        )

    async def add_user(
        self,
        user_id: int,
    ) -> bool:

        async with self.users_lock:
            users = await asyncio.to_thread(
                safe_json_load,
                self.users_file,
                [],
            )

            if isinstance(
                users,
                dict,
            ):
                raw_users = users.get(
                    "users",
                    [],
                )
            else:
                raw_users = users

            if not isinstance(
                raw_users,
                list,
            ):
                raw_users = []

            normalized: list[int] = []

            for value in raw_users:
                try:
                    normalized.append(
                        int(value)
                    )
                except (
                    TypeError,
                    ValueError,
                ):
                    pass

            normalized = list(
                dict.fromkeys(
                    normalized
                )
            )

            if user_id in normalized:
                return False

            normalized.append(
                user_id
            )

            await asyncio.to_thread(
                safe_json_save,
                self.users_file,
                normalized,
            )

            return True

    async def remove_user(
        self,
        user_id: int,
    ) -> None:

        async with self.users_lock:
            users = await asyncio.to_thread(
                safe_json_load,
                self.users_file,
                [],
            )

            if isinstance(
                users,
                dict,
            ):
                raw_users = users.get(
                    "users",
                    [],
                )
            else:
                raw_users = users

            if not isinstance(
                raw_users,
                list,
            ):
                return

            result: list[int] = []

            for value in raw_users:
                try:
                    current = int(value)
                except (
                    TypeError,
                    ValueError,
                ):
                    continue

                if current != user_id:
                    result.append(
                        current
                    )

            await asyncio.to_thread(
                safe_json_save,
                self.users_file,
                result,
            )

    # -------------------------------------------------------------------------
    # SUBSCRIPTIONS / REFERRALS / LIMITS
    # -------------------------------------------------------------------------

    @staticmethod
    def default_profile() -> dict[str, Any]:
        return {
            "subscription_until": 0.0,
            "was_referred": False,
            "referred_by": None,
            "referral_count": 0,
            "limit_used": 0,
            "limit_reset_at": 0.0,
        }

    async def get_profile(
        self,
        user_id: int,
    ) -> dict[str, Any]:

        async with self.subscription_lock:
            data = await asyncio.to_thread(
                safe_json_load,
                self.subscriptions_file,
                {},
            )

        if not isinstance(
            data,
            dict,
        ):
            data = {}

        raw_profile = data.get(
            str(user_id)
        )

        profile = self.default_profile()

        if isinstance(
            raw_profile,
            dict,
        ):
            profile.update(
                raw_profile
            )

        return profile

    async def ensure_profile(
        self,
        user_id: int,
    ) -> dict[str, Any]:

        async with self.subscription_lock:
            data = await asyncio.to_thread(
                safe_json_load,
                self.subscriptions_file,
                {},
            )

            if not isinstance(
                data,
                dict,
            ):
                data = {}

            key = str(user_id)

            profile = self.default_profile()

            raw_profile = data.get(
                key
            )

            if isinstance(
                raw_profile,
                dict,
            ):
                profile.update(
                    raw_profile
                )

            data[key] = profile

            await asyncio.to_thread(
                safe_json_save,
                self.subscriptions_file,
                data,
            )

            return profile

    async def process_referral(
        self,
        new_user_id: int,
        referrer_id: int,
    ) -> tuple[bool, int | None]:

        if new_user_id == referrer_id:
            return False, None

        async with self.users_lock:
            users = await asyncio.to_thread(
                safe_json_load,
                self.users_file,
                [],
            )

            if isinstance(
                users,
                dict,
            ):
                raw_users = users.get(
                    "users",
                    [],
                )
            else:
                raw_users = users

            if not isinstance(
                raw_users,
                list,
            ):
                raw_users = []

            normalized_users: list[int] = []

            for value in raw_users:
                try:
                    normalized_users.append(
                        int(value)
                    )
                except (
                    TypeError,
                    ValueError,
                ):
                    continue

            normalized_users = list(
                dict.fromkeys(
                    normalized_users
                )
            )

            # Новый пользователь должен быть действительно новым.
            if new_user_id in normalized_users:
                return False, None

            # Пригласивший должен существовать.
            if referrer_id not in normalized_users:
                return False, None

            normalized_users.append(
                new_user_id
            )

            await asyncio.to_thread(
                safe_json_save,
                self.users_file,
                normalized_users,
            )

            async with self.subscription_lock:
                data = await asyncio.to_thread(
                    safe_json_load,
                    self.subscriptions_file,
                    {},
                )

                if not isinstance(
                    data,
                    dict,
                ):
                    data = {}

                new_key = str(
                    new_user_id
                )

                ref_key = str(
                    referrer_id
                )

                new_profile = (
                    self.default_profile()
                )

                ref_profile = (
                    self.default_profile()
                )

                existing_new = data.get(
                    new_key
                )

                existing_ref = data.get(
                    ref_key
                )

                if isinstance(
                    existing_new,
                    dict,
                ):
                    new_profile.update(
                        existing_new
                    )

                if isinstance(
                    existing_ref,
                    dict,
                ):
                    ref_profile.update(
                        existing_ref
                    )

                if bool(
                    new_profile.get(
                        "was_referred",
                        False,
                    )
                ):
                    return False, None

                new_profile[
                    "was_referred"
                ] = True

                new_profile[
                    "referred_by"
                ] = referrer_id

                try:
                    current_count = int(
                        ref_profile.get(
                            "referral_count",
                            0,
                        )
                    )
                except (
                    TypeError,
                    ValueError,
                ):
                    current_count = 0

                ref_profile[
                    "referral_count"
                ] = current_count + 1

                now = utc_now_timestamp()

                try:
                    current_expiry = float(
                        ref_profile.get(
                            "subscription_until",
                            0.0,
                        )
                    )
                except (
                    TypeError,
                    ValueError,
                ):
                    current_expiry = 0.0

                base_time = max(
                    now,
                    current_expiry,
                )

                ref_profile[
                    "subscription_until"
                ] = (
                    base_time
                    + REFERRAL_BONUS_DAYS
                    * 24
                    * 60
                    * 60
                )

                data[new_key] = (
                    new_profile
                )

                data[ref_key] = (
                    ref_profile
                )

                await asyncio.to_thread(
                    safe_json_save,
                    self.subscriptions_file,
                    data,
                )

            return True, referrer_id

    async def add_subscription_days(
        self,
        user_id: int,
        days: int,
    ) -> float:

        async with self.subscription_lock:
            data = await asyncio.to_thread(
                safe_json_load,
                self.subscriptions_file,
                {},
            )

            if not isinstance(
                data,
                dict,
            ):
                data = {}

            key = str(user_id)

            profile = (
                self.default_profile()
            )

            existing = data.get(
                key
            )

            if isinstance(
                existing,
                dict,
            ):
                profile.update(
                    existing
                )

            now = utc_now_timestamp()

            try:
                current_expiry = float(
                    profile.get(
                        "subscription_until",
                        0.0,
                    )
                )
            except (
                TypeError,
                ValueError,
            ):
                current_expiry = 0.0

            new_expiry = (
                max(
                    now,
                    current_expiry,
                )
                + days * 24 * 60 * 60
            )

            profile[
                "subscription_until"
            ] = new_expiry

            data[key] = profile

            await asyncio.to_thread(
                safe_json_save,
                self.subscriptions_file,
                data,
            )

            return new_expiry

    async def consume_limit(
        self,
        user_id: int,
    ) -> tuple[bool, int, int, float]:
        """
        Списывает ровно ОДИН лимит за принятое
        пользовательское сообщение.

        Это НЕ связано с количеством запросов к AI API.

        Один принятый запрос:
            user message -> consume_limit() один раз.

        Скрытое сжатие истории лимит не списывает.

        Если limit = 20:
            used 0..19 -> разрешается;
            used 20      -> блокируется.

        Операция атомарна под subscription_lock.
        """

        async with self.subscription_lock:
            data = await asyncio.to_thread(
                safe_json_load,
                self.subscriptions_file,
                {},
            )

            if not isinstance(
                data,
                dict,
            ):
                data = {}

            key = str(
                user_id
            )

            profile = (
                self.default_profile()
            )

            existing = data.get(
                key
            )

            if isinstance(
                existing,
                dict,
            ):
                profile.update(
                    existing
                )

            now = utc_now_timestamp()

            try:
                subscription_until = float(
                    profile.get(
                        "subscription_until",
                        0.0,
                    )
                )
            except (
                TypeError,
                ValueError,
            ):
                subscription_until = 0.0

            subscribed = (
                subscription_until > now
            )

            limit = (
                SUBSCRIPTION_LIMIT
                if subscribed
                else FREE_LIMIT
            )

            try:
                reset_at = float(
                    profile.get(
                        "limit_reset_at",
                        0.0,
                    )
                )
            except (
                TypeError,
                ValueError,
            ):
                reset_at = 0.0

            try:
                used = int(
                    profile.get(
                        "limit_used",
                        0,
                    )
                )
            except (
                TypeError,
                ValueError,
            ):
                used = 0

            # -----------------------------------------------------------------
            # Новый период.
            # -----------------------------------------------------------------

            if (
                LIMIT_PERIOD_SECONDS > 0
                and (
                    reset_at <= 0
                    or now >= reset_at
                )
            ):
                used = 0

                reset_at = (
                    now
                    + LIMIT_PERIOD_SECONDS
                )

            # -----------------------------------------------------------------
            # Граница лимита.
            # -----------------------------------------------------------------

            if used >= limit:
                profile[
                    "limit_used"
                ] = used

                profile[
                    "limit_reset_at"
                ] = reset_at

                data[key] = profile

                await asyncio.to_thread(
                    safe_json_save,
                    self.subscriptions_file,
                    data,
                )

                return (
                    False,
                    used,
                    limit,
                    reset_at,
                )

            # Ровно один пользовательский запрос.
            used += 1

            profile[
                "limit_used"
            ] = used

            profile[
                "limit_reset_at"
            ] = reset_at

            data[key] = profile

            await asyncio.to_thread(
                safe_json_save,
                self.subscriptions_file,
                data,
            )

            return (
                True,
                used,
                limit,
                reset_at,
            )

    # -------------------------------------------------------------------------
    # HISTORY
    # -------------------------------------------------------------------------

    def history_path(
        self,
        user_id: int,
    ) -> Path:
        return (
            self.history_dir
            / f"{user_id}.json"
        )

    async def get_history(
        self,
        user_id: int,
    ) -> list[dict[str, Any]]:

        path = self.history_path(
            user_id
        )

        data = await asyncio.to_thread(
            safe_json_load,
            path,
            [],
        )

        if not isinstance(
            data,
            list,
        ):
            return []

        result: list[dict[str, Any]] = []

        for item in data:
            if not isinstance(
                item,
                dict,
            ):
                continue

            role = item.get(
                "role"
            )

            content = item.get(
                "content"
            )

            if role in {
                "system",
                "user",
                "assistant",
            }:
                result.append(
                    {
                        "role": role,
                        "content": content,
                    }
                )

        return result

    async def save_history(
        self,
        user_id: int,
        history: list[dict[str, Any]],
    ) -> None:

        path = self.history_path(
            user_id
        )

        await asyncio.to_thread(
            safe_json_save,
            path,
            history,
        )


# =============================================================================
# FILE PROCESSOR
# =============================================================================

class FileProcessor:
    def __init__(self):
        self.semaphore = asyncio.Semaphore(
            MAX_CONCURRENT_FILE_PROCESSING
        )

    async def download_file(
        self,
        bot: Bot,
        telegram_file_id: str,
        original_name: str,
    ) -> Path:

        extension = get_file_extension(
            original_name
        )

        if extension not in SUPPORTED_EXTENSIONS:
            raise ValueError(
                "Этот тип файла не поддерживается."
            )

        path = (
            FILES_DIR
            / f"{uuid.uuid4().hex}{extension}"
        )

        telegram_file = await bot.get_file(
            telegram_file_id
        )

        if not telegram_file.file_path:
            raise RuntimeError(
                "Telegram не вернул путь к файлу."
            )

        await bot.download_file(
            telegram_file.file_path,
            destination=path,
        )

        try:
            size = path.stat().st_size
        except FileNotFoundError:
            raise RuntimeError(
                "Файл не был сохранён."
            )

        if size > MAX_FILE_SIZE:
            try:
                path.unlink()
            except Exception:
                pass

            raise ValueError(
                "Файл слишком большой. "
                f"Максимальный размер: "
                f"{MAX_FILE_SIZE_MB} МБ."
            )

        return path

    async def process_path(
        self,
        path: Path,
        original_name: str,
        mime_type: str | None,
    ) -> dict[str, Any]:

        async with self.semaphore:
            try:
                return await asyncio.wait_for(
                    asyncio.to_thread(
                        self._process_sync,
                        path,
                        original_name,
                        mime_type,
                    ),
                    timeout=AI_TIMEOUT_SECONDS,
                )

            except asyncio.TimeoutError:
                raise RuntimeError(
                    "Обработка файла заняла слишком много времени."
                )

    def _process_sync(
        self,
        path: Path,
        original_name: str,
        mime_type: str | None,
    ) -> dict[str, Any]:

        extension = path.suffix.lower()

        if extension == ".pdf":
            reader = PdfReader(
                str(path)
            )

            pages: list[str] = []

            for page in reader.pages:
                try:
                    text = (
                        page.extract_text()
                        or ""
                    )
                except Exception:
                    text = ""

                if text.strip():
                    pages.append(
                        text.strip()
                    )

            return {
                "type": "text",
                "name": original_name,
                "text": "\n\n".join(
                    pages
                ).strip(),
            }

        if extension == ".docx":
            document = Document(
                str(path)
            )

            parts: list[str] = []

            for paragraph in document.paragraphs:
                text = paragraph.text.strip()

                if text:
                    parts.append(
                        text
                    )

            for table in document.tables:
                for row in table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                    ]

                    if any(cells):
                        parts.append(
                            " | ".join(cells)
                        )

            return {
                "type": "text",
                "name": original_name,
                "text": "\n".join(
                    parts
                ).strip(),
            }

        if extension == ".xlsx":
            workbook = load_workbook(
                filename=str(path),
                read_only=True,
                data_only=True,
            )

            parts: list[str] = []

            try:
                for sheet in workbook.worksheets:
                    parts.append(
                        f"Лист: {sheet.title}"
                    )

                    for row in sheet.iter_rows(
                        values_only=True
                    ):
                        values = []

                        for value in row:
                            if value is None:
                                values.append("")
                            else:
                                values.append(
                                    str(value)
                                )

                        if any(
                            value.strip()
                            for value in values
                        ):
                            parts.append(
                                " | ".join(
                                    values
                                )
                            )

            finally:
                workbook.close()

            return {
                "type": "text",
                "name": original_name,
                "text": "\n".join(
                    parts
                ).strip(),
            }

        if extension == ".pptx":
            presentation = Presentation(
                str(path)
            )

            parts: list[str] = []

            for index, slide in enumerate(
                presentation.slides,
                start=1,
            ):
                parts.append(
                    f"Слайд {index}:"
                )

                for shape in slide.shapes:
                    if not hasattr(
                        shape,
                        "text",
                    ):
                        continue

                    text = shape.text.strip()

                    if text:
                        parts.append(
                            text
                        )

            return {
                "type": "text",
                "name": original_name,
                "text": "\n".join(
                    parts
                ).strip(),
            }

        if extension in {
            ".jpg",
            ".jpeg",
            ".png",
            ".webp",
        }:
            if extension in {
                ".jpg",
                ".jpeg",
            }:
                detected_mime = (
                    "image/jpeg"
                )
            elif extension == ".png":
                detected_mime = (
                    "image/png"
                )
            else:
                detected_mime = (
                    "image/webp"
                )

            return {
                "type": "image",
                "name": original_name,
                "mime_type": detected_mime,
                "data_url": image_to_data_url(
                    path,
                    detected_mime,
                ),
            }

        raise ValueError(
            "Этот тип файла не поддерживается."
        )


# =============================================================================
# AI SERVICE
# =============================================================================

class AIService:
    def __init__(self):
        self.session: ClientSession | None = None

        self.semaphore = asyncio.Semaphore(
            MAX_CONCURRENT_AI_REQUESTS
        )

    async def start(self) -> None:
        if self.session is not None:
            return

        timeout = ClientTimeout(
            total=AI_TIMEOUT_SECONDS,
            connect=min(
                20,
                AI_TIMEOUT_SECONDS,
            ),
        )

        self.session = ClientSession(
            timeout=timeout,
            headers={
                "Authorization": (
                    f"Bearer {AI_API_KEY}"
                ),
                "Content-Type": (
                    "application/json"
                ),
            },
        )

    async def close(self) -> None:
        if self.session is not None:
            await self.session.close()
            self.session = None

    def get_system_prompt(self) -> str:
        return (
            f"Ты — {AI_MODEL_NAME}, "
            "AI-модель, работающая внутри Telegram-бота.\n\n"
            f"{SYSTEM_PROMPT}"
        )

    @staticmethod
    def normalize_messages(
        messages: list[dict[str, Any]],
    ) -> list[dict[str, Any]]:

        system_messages = [
            message
            for message in messages
            if message.get("role") == "system"
        ]

        other_messages = [
            message
            for message in messages
            if message.get("role") != "system"
        ]

        if system_messages:
            return [
                system_messages[0],
                *other_messages,
            ]

        return other_messages

    async def chat(
        self,
        messages: list[dict[str, Any]],
        *,
        temperature: float | None = None,
    ) -> str:

        if self.session is None:
            raise RuntimeError(
                "AI service не запущен."
            )

        normalized = (
            self.normalize_messages(
                messages
            )
        )

        payload: dict[str, Any] = {
            "model": AI_MODEL_ID,
            "messages": normalized,
        }

        if temperature is not None:
            payload["temperature"] = (
                temperature
            )

        url = (
            f"{AI_BASE_URL}"
            "/chat/completions"
        )

        async with self.semaphore:
            try:
                async with self.session.post(
                    url,
                    json=payload,
                ) as response:

                    raw_text = (
                        await response.text()
                    )

                    if response.status >= 400:
                        raise RuntimeError(
                            "AI API error "
                            f"{response.status}: "
                            f"{raw_text[:4000]}"
                        )

                    try:
                        data = json.loads(
                            raw_text
                        )
                    except json.JSONDecodeError:
                        raise RuntimeError(
                            "AI API вернул "
                            "некорректный JSON."
                        )

                    return self.extract_content(
                        data
                    )

            except asyncio.TimeoutError:
                raise RuntimeError(
                    "AI API не ответил вовремя."
                )

    @staticmethod
    def extract_content(
        data: dict[str, Any],
    ) -> str:

        choices = data.get(
            "choices"
        )

        if (
            not isinstance(
                choices,
                list,
            )
            or not choices
        ):
            raise RuntimeError(
                "AI API не вернул choices."
            )

        choice = choices[0]

        if not isinstance(
            choice,
            dict,
        ):
            raise RuntimeError(
                "AI API вернул "
                "некорректный choice."
            )

        message = choice.get(
            "message"
        )

        if not isinstance(
            message,
            dict,
        ):
            raise RuntimeError(
                "AI API не вернул message."
            )

        content = message.get(
            "content"
        )

        if isinstance(
            content,
            str,
        ):
            return content.strip()

        if isinstance(
            content,
            list,
        ):
            parts: list[str] = []

            for part in content:
                if not isinstance(
                    part,
                    dict,
                ):
                    continue

                text = part.get(
                    "text"
                )

                if isinstance(
                    text,
                    str,
                ):
                    parts.append(
                        text
                    )

            result = "\n".join(
                parts
            ).strip()

            if result:
                return result

        raise RuntimeError(
            "AI API не вернул "
            "текст ответа."
        )


# =============================================================================
# QUEUE
# =============================================================================

class QueueFullError(Exception):
    pass


class QueueItem:
    def __init__(
        self,
        future: asyncio.Future,
        messages: list[dict[str, Any]],
    ):
        self.future = future
        self.messages = messages


class AIQueue:
    def __init__(
        self,
        ai: AIService,
    ):
        self.ai = ai

        self.queue: asyncio.Queue[
            QueueItem
        ] = asyncio.Queue()

        self.worker_tasks: list[
            asyncio.Task
        ] = []

        self.running = False

    async def start(self) -> None:
        if self.running:
            return

        self.running = True

        self.worker_tasks = [
            asyncio.create_task(
                self.worker_loop(),
                name=(
                    f"ai-worker-{index + 1}"
                ),
            )
            for index in range(
                MAX_CONCURRENT_AI_REQUESTS
            )
        ]

    async def stop(self) -> None:
        self.running = False

        tasks = list(
            self.worker_tasks
        )

        self.worker_tasks.clear()

        for task in tasks:
            task.cancel()

        for task in tasks:
            try:
                await task
            except asyncio.CancelledError:
                pass

    async def submit(
        self,
        messages: list[dict[str, Any]],
    ) -> str:

        if not self.running:
            raise RuntimeError(
                "Очередь AI не запущена."
            )

        loop = (
            asyncio.get_running_loop()
        )

        future: asyncio.Future = (
            loop.create_future()
        )

        item = QueueItem(
            future=future,
            messages=messages,
        )

        try:
            self.queue.put_nowait(
                item
            )
        except asyncio.QueueFull:
            raise QueueFullError

        try:
            return await future

        except asyncio.CancelledError:
            if not future.done():
                future.cancel()

            raise

    async def worker_loop(self) -> None:
        while self.running:
            item = await self.queue.get()

            try:
                if item.future.cancelled():
                    continue

                try:
                    result = await self.ai.chat(
                        item.messages
                    )

                    if not item.future.done():
                        item.future.set_result(
                            result
                        )

                except asyncio.CancelledError:
                    if not item.future.done():
                        item.future.cancel()

                    raise

                except Exception as exc:
                    if not item.future.done():
                        item.future.set_exception(
                            exc
                        )

            finally:
                self.queue.task_done()


# =============================================================================
# BOT APPLICATION
# =============================================================================

class BotApp:
    def __init__(self):
        self.bot = Bot(
            BOT_TOKEN,
            default=DefaultBotProperties(
                parse_mode=ParseMode.HTML
            ),
        )

        self.dp = Dispatcher()
        self.router = Router()

        self.store = JSONStore()
        self.files = FileProcessor()
        self.ai = AIService()
        self.queue = AIQueue(
            self.ai
        )

        self.active_tasks: dict[
            int,
            asyncio.Task,
        ] = {}

        self.active_lock = asyncio.Lock()

        self.broadcast_lock = (
            asyncio.Lock()
        )

        self.bot_username = ""

        self.register_handlers()

    # =========================================================================
    # HANDLERS
    # =========================================================================

    def register_handlers(self) -> None:
        self.router.message.register(
            self.handle_start,
            CommandStart(),
        )

        self.router.message.register(
            self.handle_subscription_command,
            Command(
                SUBSCRIPTION_COMMAND
            ),
        )

        self.router.message.register(
            self.handle_subscription_button,
            F.text == SUBSCRIPTION_BUTTON,
        )

        self.router.message.register(
            self.handle_send,
            Command("send"),
        )

        self.router.callback_query.register(
            self.handle_stop,
            F.data == STOP_CALLBACK,
        )

        self.router.message.register(
            self.handle_document,
            F.document,
        )

        self.router.message.register(
            self.handle_photo,
            F.photo,
        )

        self.router.message.register(
            self.handle_text,
            F.text,
        )

        self.router.message.register(
            self.handle_unsupported,
        )

        self.dp.include_router(
            self.router
        )

    # =========================================================================
    # LIFECYCLE
    # =========================================================================

    async def start(self) -> None:
        if not BOT_TOKEN:
            raise RuntimeError(
                "BOT_TOKEN не задан"
            )

        if not AI_BASE_URL:
            raise RuntimeError(
                "AI_BASE_URL не задан"
            )

        if not AI_API_KEY:
            raise RuntimeError(
                "AI_API_KEY не задан"
            )

        if not AI_MODEL_ID:
            raise RuntimeError(
                "AI_MODEL_ID не задан"
            )

        me = await self.bot.get_me()

        if not me.username:
            raise RuntimeError(
                "Telegram не вернул username бота."
            )

        self.bot_username = me.username

        await self.ai.start()
        await self.queue.start()

        log.info(
            "Bot started | model=%s | model_id=%s | username=@%s",
            AI_MODEL_NAME,
            AI_MODEL_ID,
            self.bot_username,
        )

        log.info(
            "Limits | free=%s | subscription=%s | period=%sh | referral=%sd",
            FREE_LIMIT,
            SUBSCRIPTION_LIMIT,
            LIMIT_PERIOD_HOURS,
            REFERRAL_BONUS_DAYS,
        )

        log.info(
            "Admin IDs configured: %s",
            sorted(ADMIN_IDS),
        )

    async def stop(self) -> None:
        async with self.active_lock:
            tasks = list(
                self.active_tasks.values()
            )

        for task in tasks:
            task.cancel()

        for task in tasks:
            try:
                await task
            except asyncio.CancelledError:
                pass
            except Exception:
                pass

        await self.queue.stop()
        await self.ai.close()

        try:
            await self.bot.session.close()
        except Exception:
            pass

    # =========================================================================
    # ACTIVE REQUEST LOCK
    # =========================================================================

    async def is_busy(
        self,
        user_id: int,
    ) -> bool:

        async with self.active_lock:
            task = self.active_tasks.get(
                user_id
            )

            if task is None:
                return False

            if task.done():
                self.active_tasks.pop(
                    user_id,
                    None,
                )

                return False

            return True

    async def register_task(
        self,
        user_id: int,
        task: asyncio.Task,
    ) -> bool:

        async with self.active_lock:
            current = self.active_tasks.get(
                user_id
            )

            if (
                current is not None
                and not current.done()
            ):
                return False

            self.active_tasks[user_id] = task

            return True

    async def release_task(
        self,
        user_id: int,
        task: asyncio.Task,
    ) -> None:

        async with self.active_lock:
            current = self.active_tasks.get(
                user_id
            )

            if current is task:
                self.active_tasks.pop(
                    user_id,
                    None,
                )

    # =========================================================================
    # REPLY KEYBOARD
    # =========================================================================

    @staticmethod
    def main_keyboard() -> ReplyKeyboardMarkup:
        return ReplyKeyboardMarkup(
            keyboard=[
                [
                    KeyboardButton(
                        text=SUBSCRIPTION_BUTTON
                    )
                ],
            ],
            resize_keyboard=True,
        )

    # =========================================================================
    # STOP KEYBOARD
    # =========================================================================

    @staticmethod
    def stop_keyboard() -> InlineKeyboardMarkup:
        return InlineKeyboardMarkup(
            inline_keyboard=[
                [
                    InlineKeyboardButton(
                        text="⏹ Стоп",
                        callback_data=STOP_CALLBACK,
                    )
                ]
            ]
        )

    async def handle_stop(
        self,
        callback: CallbackQuery,
    ) -> None:

        if callback.from_user is None:
            await callback.answer()
            return

        user_id = callback.from_user.id

        async with self.active_lock:
            task = self.active_tasks.get(
                user_id
            )

        if task is None or task.done():
            await callback.answer(
                "Активного запроса нет."
            )
            return

        task.cancel()

        await callback.answer(
            "Запрос остановлен."
        )

        try:
            if callback.message:
                await callback.message.edit_reply_markup(
                    reply_markup=None
                )
        except Exception:
            pass

    # =========================================================================
    # SUBSCRIPTION
    # =========================================================================

    async def get_referral_link(
        self,
        user_id: int,
    ) -> str:

        if not self.bot_username:
            me = await self.bot.get_me()

            if not me.username:
                raise RuntimeError(
                    "Не удалось получить username бота."
                )

            self.bot_username = me.username

        return (
            f"https://t.me/"
            f"{self.bot_username}"
            f"?start=ref_{user_id}"
        )

    async def show_subscription_menu(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        user_id = message.from_user.id

        await self.store.add_user(
            user_id
        )

        profile = (
            await self.store.ensure_profile(
                user_id
            )
        )

        now = utc_now_timestamp()

        try:
            subscription_until = float(
                profile.get(
                    "subscription_until",
                    0.0,
                )
            )
        except (
            TypeError,
            ValueError,
        ):
            subscription_until = 0.0

        active = (
            subscription_until > now
        )

        try:
            referral_count = int(
                profile.get(
                    "referral_count",
                    0,
                )
            )
        except (
            TypeError,
            ValueError,
        ):
            referral_count = 0

        referral_link = (
            await self.get_referral_link(
                user_id
            )
        )

        if active:
            subscription_status = (
                "🟢 <b>Подписка активна</b>\n"
                "Осталось: <b>"
                f"{format_remaining_subscription(subscription_until)}"
                "</b>\n"
                "До: <b>"
                f"{format_datetime(subscription_until)}"
                "</b>"
            )

            current_limit = (
                SUBSCRIPTION_LIMIT
            )

        else:
            subscription_status = (
                "⚪ <b>Подписка неактивна</b>"
            )

            current_limit = FREE_LIMIT

        if SUPPORT_USERNAME:
            support_text = (
                "\n\n"
                "💬 Поддержка: "
                f"@{SUPPORT_USERNAME}"
            )

        elif ADMIN_IDS:
            support_text = (
                "\n\n"
                "💬 По вопросам покупки "
                "обратитесь к администратору."
            )

        else:
            support_text = (
                "\n\n"
                "💬 Поддержка временно "
                "не настроена."
            )

        text = (
            "👥 <b>Подписка</b>\n\n"
            f"{subscription_status}\n\n"
            "📊 Лимит запросов: "
            f"<b>{current_limit}</b> "
            f"за {LIMIT_PERIOD_HOURS} ч.\n\n"
            "🎁 <b>Пригласите друга</b>\n"
            "За каждого нового пользователя "
            "вы получите "
            f"<b>{REFERRAL_BONUS_DAYS} дн.</b> "
            "подписки.\n\n"
            "🔗 Ваша реферальная ссылка:\n"
            f"<code>{html.escape(referral_link)}</code>\n\n"
            "👤 Приглашено: "
            f"<b>{referral_count}</b>"
            f"{support_text}"
        )

        keyboard_rows: list[
            list[InlineKeyboardButton]
        ] = []

        if SUPPORT_USERNAME:
            keyboard_rows.append(
                [
                    InlineKeyboardButton(
                        text="💬 Поддержка",
                        url=(
                            f"https://t.me/"
                            f"{SUPPORT_USERNAME}"
                        ),
                    )
                ]
            )

        await message.answer(
            text,
            reply_markup=(
                InlineKeyboardMarkup(
                    inline_keyboard=keyboard_rows
                )
                if keyboard_rows
                else None
            ),
        )

    async def handle_subscription_command(
        self,
        message: Message,
    ) -> None:
        await self.show_subscription_menu(
            message
        )

    async def handle_subscription_button(
        self,
        message: Message,
    ) -> None:
        await self.show_subscription_menu(
            message
        )

    # =========================================================================
    # START
    # =========================================================================

    async def handle_start(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        user_id = message.from_user.id

        text = message.text or ""

        parts = text.split(
            maxsplit=1
        )

        referral_id: int | None = None

        if len(parts) == 2:
            payload = parts[1].strip()

            if payload.startswith("ref_"):
                raw_referrer = (
                    payload[4:].strip()
                )

                if raw_referrer.isdigit():
                    try:
                        referral_id = int(
                            raw_referrer
                        )
                    except ValueError:
                        referral_id = None

        referral_success = False

        if referral_id is not None:
            try:
                (
                    referral_success,
                    _,
                ) = await self.store.process_referral(
                    user_id,
                    referral_id,
                )

            except Exception:
                log.exception(
                    "Referral processing failed | "
                    "user=%s referrer=%s",
                    user_id,
                    referral_id,
                )

        else:
            await self.store.add_user(
                user_id
            )

            await self.store.ensure_profile(
                user_id
            )

        if referral_success:
            await message.answer(
                "🎉 <b>Реферал засчитан!</b>\n\n"
                "Пригласивший получил "
                f"<b>{REFERRAL_BONUS_DAYS} дн.</b> "
                "подписки.",
                reply_markup=self.main_keyboard(),
            )

        else:
            await message.answer(
                "👋 Привет!\n\n"
                f"Я — {AI_MODEL_NAME}.\n"
                "Отправьте мне сообщение "
                "или файл.\n\n"
            )

    # =========================================================================
    # ADMIN BROADCAST
    # =========================================================================

    async def handle_send(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        user_id = message.from_user.id

        if user_id not in ADMIN_IDS:
            await message.answer(
                "⛔ Доступ запрещён."
            )
            return

        text = message.text or ""

        parts = text.split(
            maxsplit=1
        )

        if len(parts) < 2:
            await message.answer(
                "Использование:\n"
                "<code>/send текст рассылки</code>"
            )
            return

        broadcast_text = parts[1].strip()

        if not broadcast_text:
            await message.answer(
                "Текст рассылки пустой."
            )
            return

        asyncio.create_task(
            self.broadcast(
                broadcast_text,
                user_id,
            )
        )

        await message.answer(
            "📨 Рассылка запущена."
        )

    async def broadcast(
        self,
        text: str,
        admin_id: int,
    ) -> None:

        async with self.broadcast_lock:
            users = await self.store.get_users()

            total = len(users)

            if total == 0:
                try:
                    await self.bot.send_message(
                        admin_id,
                        "📨 Пользователей "
                        "для рассылки нет.",
                    )
                except Exception:
                    pass

                return

            try:
                status_message = (
                    await self.bot.send_message(
                        admin_id,
                        "📨 Рассылка запущена.\n"
                        f"Получателей: {total}",
                    )
                )
            except Exception:
                status_message = None

            semaphore = asyncio.Semaphore(
                BROADCAST_CONCURRENCY
            )

            sent = 0
            failed = 0
            removed = 0

            counter_lock = asyncio.Lock()

            async def send_one(
                target_id: int,
            ) -> None:

                nonlocal sent
                nonlocal failed
                nonlocal removed

                async with semaphore:
                    try:
                        await self.bot.send_message(
                            target_id,
                            text,
                        )

                        async with counter_lock:
                            sent += 1

                    except TelegramRetryAfter as exc:
                        try:
                            await asyncio.sleep(
                                exc.retry_after
                            )

                            await self.bot.send_message(
                                target_id,
                                text,
                            )

                            async with counter_lock:
                                sent += 1

                        except Exception:
                            async with counter_lock:
                                failed += 1

                    except (
                        TelegramForbiddenError,
                        TelegramBadRequest,
                    ):
                        await self.store.remove_user(
                            target_id
                        )

                        async with counter_lock:
                            removed += 1

                    except TelegramNetworkError:
                        async with counter_lock:
                            failed += 1

                    except Exception:
                        log.exception(
                            "Broadcast failed for %s",
                            target_id,
                        )

                        async with counter_lock:
                            failed += 1

                    finally:
                        if BROADCAST_DELAY > 0:
                            await asyncio.sleep(
                                BROADCAST_DELAY
                            )

            batch_size = max(
                1,
                BROADCAST_CONCURRENCY * 10,
            )

            for start in range(
                0,
                len(users),
                batch_size,
            ):
                batch = users[
                    start:start + batch_size
                ]

                await asyncio.gather(
                    *(
                        send_one(target_id)
                        for target_id in batch
                    )
                )

            final_text = (
                "📨 Рассылка завершена.\n\n"
                f"Всего: {total}\n"
                f"Отправлено: {sent}\n"
                f"Ошибок: {failed}\n"
                "Удалено заблокированных: "
                f"{removed}"
            )

            if status_message is not None:
                try:
                    await status_message.edit_text(
                        final_text
                    )
                    return
                except Exception:
                    pass

            try:
                await self.bot.send_message(
                    admin_id,
                    final_text,
                )
            except Exception:
                pass

    # =========================================================================
    # COMMON REQUEST ADMISSION
    # =========================================================================

    async def begin_user_request(
        self,
        message: Message,
    ) -> Optional[asyncio.Task]:

        if not message.from_user:
            return None

        user_id = message.from_user.id

        # ---------------------------------------------------------------------
        # На одного пользователя одновременно только один принятый запрос.
        # ---------------------------------------------------------------------

        if await self.is_busy(user_id):
            await message.answer(
                BUSY_TEXT
            )
            return None

        # ---------------------------------------------------------------------
        # Task регистрируется сразу.
        # ---------------------------------------------------------------------

        task = asyncio.create_task(
            self.process_request_wrapper(
                message
            ),
            name=(
                f"user-request-{user_id}"
            ),
        )

        registered = await self.register_task(
            user_id,
            task,
        )

        if not registered:
            task.cancel()

            try:
                await task
            except asyncio.CancelledError:
                pass
            except Exception:
                pass

            await message.answer(
                BUSY_TEXT
            )

            return None

        return task

    # =========================================================================
    # TEXT
    # =========================================================================

    async def handle_text(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        text = (
            message.text or ""
        ).strip()

        if not text:
            return

        await self.store.add_user(
            message.from_user.id
        )

        await self.store.ensure_profile(
            message.from_user.id
        )

        await self.begin_user_request(
            message
        )

    # =========================================================================
    # DOCUMENT
    # =========================================================================

    async def handle_document(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        await self.store.add_user(
            message.from_user.id
        )

        await self.store.ensure_profile(
            message.from_user.id
        )

        document = message.document

        if document is None:
            return

        file_name = (
            document.file_name
            or "file"
        )

        extension = get_file_extension(
            file_name
        )

        if extension not in SUPPORTED_EXTENSIONS:
            await message.answer(
                "❌ Этот тип файла "
                "не поддерживается.\n\n"
                "Поддерживаются: PDF, DOCX, "
                "XLSX, PPTX, JPG, JPEG, PNG "
                "и WEBP."
            )
            return

        if (
            document.file_size is not None
            and document.file_size > MAX_FILE_SIZE
        ):
            await message.answer(
                "❌ Файл слишком большой.\n"
                "Максимальный размер: "
                f"{MAX_FILE_SIZE_MB} МБ."
            )
            return

        await self.begin_user_request(
            message
        )

    # =========================================================================
    # PHOTO
    # =========================================================================

    async def handle_photo(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        await self.store.add_user(
            message.from_user.id
        )

        await self.store.ensure_profile(
            message.from_user.id
        )

        if not message.photo:
            return

        largest = message.photo[-1]

        if (
            largest.file_size is not None
            and largest.file_size > MAX_FILE_SIZE
        ):
            await message.answer(
                "❌ Изображение слишком большое.\n"
                "Максимальный размер: "
                f"{MAX_FILE_SIZE_MB} МБ."
            )
            return

        await self.begin_user_request(
            message
        )

    # =========================================================================
    # UNSUPPORTED
    # =========================================================================

    async def handle_unsupported(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        await message.answer(
            "❌ Этот тип сообщения "
            "не поддерживается.\n\n"
            "Можно отправить текст, PDF, DOCX, "
            "XLSX, PPTX, JPG, JPEG, PNG или WEBP."
        )

    # =========================================================================
    # REQUEST WRAPPER
    # =========================================================================

    async def process_request_wrapper(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        user_id = message.from_user.id

        current_task = (
            asyncio.current_task()
        )

        try:
            await self.process_request(
                message
            )

        except asyncio.CancelledError:
            try:
                await message.answer(
                    STOPPED_TEXT
                )
            except Exception:
                pass

            raise

        except QueueFullError:
            try:
                await message.answer(
                    QUEUE_FULL_TEXT
                )
            except Exception:
                pass

        except Exception:
            log.exception(
                "Request failed for user %s",
                user_id,
            )

            try:
                await message.answer(
                    "❌ Произошла ошибка "
                    "при обработке запроса."
                )
            except Exception:
                pass

        finally:
            if current_task is not None:
                await self.release_task(
                    user_id,
                    current_task,
                )

    # =========================================================================
    # LIMIT MESSAGE
    # =========================================================================

    async def send_limit_reached(
        self,
        message: Message,
        used: int,
        limit: int,
        reset_at: float,
    ) -> None:

        now = utc_now_timestamp()

        remaining = max(
            0,
            int(
                reset_at - now
            ),
        )

        hours = remaining // 3600

        minutes = (
            remaining % 3600
        ) // 60

        if hours > 0:
            reset_text = (
                f"{hours} ч."
            )

        elif minutes > 0:
            reset_text = (
                f"{minutes} мин."
            )

        else:
            reset_text = (
                "менее минуты"
            )

        profile = (
            await self.store.get_profile(
                message.from_user.id
            )
        )

        try:
            subscription_until = float(
                profile.get(
                    "subscription_until",
                    0.0,
                )
            )
        except (
            TypeError,
            ValueError,
        ):
            subscription_until = 0.0

        subscribed = (
            subscription_until > now
        )

        if subscribed:
            text = (
                "⛔ <b>Лимит запросов исчерпан.</b>\n\n"
                f"Использовано: <b>{used}/{limit}</b>\n"
                "Новый период через: "
                f"<b>{reset_text}</b>."
            )

        else:
            text = (
                "⛔ <b>Бесплатный лимит запросов исчерпан.</b>\n\n"
                f"Использовано: <b>{used}/{limit}</b>\n"
                "Новый период через: "
                f"<b>{reset_text}</b>.\n\n"
                "👥 Откройте раздел подписки, "
                "чтобы узнать условия её получения."
            )

        await message.answer(
            text
        )

    # =========================================================================
    # MAIN REQUEST
    # =========================================================================

    async def process_request(
        self,
        message: Message,
    ) -> None:

        if not message.from_user:
            return

        user_id = message.from_user.id

        # ---------------------------------------------------------------------
        # КРИТИЧЕСКАЯ ПРОВЕРКА ЛИМИТА.
        #
        # Этот вызов происходит один раз на принятое пользовательское
        # сообщение.
        #
        # Дальнейшие queue.submit(), включая скрытое сжатие истории,
        # НЕ вызывают consume_limit().
        # ---------------------------------------------------------------------

        (
            allowed,
            used,
            limit,
            reset_at,
        ) = await self.store.consume_limit(
            user_id
        )

        if not allowed:
            await self.send_limit_reached(
                message,
                used,
                limit,
                reset_at,
            )
            return

        user_text = (
            message.text or ""
        ).strip()

        files_context: list[
            dict[str, Any]
        ] = []

        # ---------------------------------------------------------------------
        # DOCUMENT
        # ---------------------------------------------------------------------

        if message.document is not None:
            document = message.document

            file_name = (
                document.file_name
                or "file"
            )

            extension = get_file_extension(
                file_name
            )

            if extension not in SUPPORTED_EXTENSIONS:
                raise ValueError(
                    "Этот тип файла "
                    "не поддерживается."
                )

            temp_path: Path | None = None

            try:
                temp_path = (
                    await self.files.download_file(
                        self.bot,
                        document.file_id,
                        file_name,
                    )
                )

                processed = (
                    await self.files.process_path(
                        temp_path,
                        file_name,
                        document.mime_type,
                    )
                )

                files_context.append(
                    processed
                )

            finally:
                if temp_path is not None:
                    try:
                        temp_path.unlink(
                            missing_ok=True
                        )
                    except Exception:
                        pass

        # ---------------------------------------------------------------------
        # PHOTO
        # ---------------------------------------------------------------------

        elif message.photo:
            photo = message.photo[-1]

            temp_path: Path | None = None

            try:
                temp_path = (
                    await self.files.download_file(
                        self.bot,
                        photo.file_id,
                        "image.jpg",
                    )
                )

                processed = (
                    await self.files.process_path(
                        temp_path,
                        "image.jpg",
                        "image/jpeg",
                    )
                )

                files_context.append(
                    processed
                )

            finally:
                if temp_path is not None:
                    try:
                        temp_path.unlink(
                            missing_ok=True
                        )
                    except Exception:
                        pass

        # ---------------------------------------------------------------------
        # HISTORY
        # ---------------------------------------------------------------------

        history = (
            await self.store.get_history(
                user_id
            )
        )

        if len(history) > MAX_HISTORY_MESSAGES:
            history = history[
                -MAX_HISTORY_MESSAGES:
            ]

        # ---------------------------------------------------------------------
        # STATUS
        # ---------------------------------------------------------------------

        status_message = await message.answer(
            "🤔 Обрабатываю запрос...",
            reply_markup=self.stop_keyboard(),
        )

        # ---------------------------------------------------------------------
        # BUILD USER MESSAGE
        # ---------------------------------------------------------------------

        text_parts: list[str] = []

        if user_text:
            text_parts.append(
                user_text
            )

        for file_data in files_context:
            if (
                file_data.get("type")
                == "text"
            ):
                file_name = (
                    file_data.get(
                        "name",
                        "файл",
                    )
                )

                file_text = (
                    file_data.get(
                        "text",
                        "",
                    )
                )

                if file_text:
                    text_parts.append(
                        "\n\n"
                        f"--- Файл: {file_name} ---\n"
                        f"{file_text}\n"
                        "--- Конец файла ---"
                    )

                else:
                    text_parts.append(
                        "\n\n"
                        f"--- Файл: {file_name} ---\n"
                        "[В файле не удалось "
                        "извлечь текст]\n"
                        "--- Конец файла ---"
                    )

        text_content = "\n".join(
            text_parts
        ).strip()

        image_parts: list[
            dict[str, Any]
        ] = []

        for file_data in files_context:
            if (
                file_data.get("type")
                != "image"
            ):
                continue

            data_url = file_data.get(
                "data_url"
            )

            if not isinstance(
                data_url,
                str,
            ):
                continue

            image_parts.append(
                {
                    "type": "image_url",
                    "image_url": {
                        "url": data_url
                    },
                }
            )

        if image_parts:
            content_parts: list[
                dict[str, Any]
            ] = []

            if text_content:
                content_parts.append(
                    {
                        "type": "text",
                        "text": text_content,
                    }
                )

            content_parts.extend(
                image_parts
            )

            main_user_content: Any = (
                content_parts
            )

        else:
            main_user_content = (
                text_content
                or "Проанализируй "
                "прикреплённый файл."
            )

        # ---------------------------------------------------------------------
        # MAIN AI REQUEST
        #
        # ВАЖНО:
        # Скрытого STAGE 1 / PLAN больше НЕТ.
        #
        # Один принятый пользовательский запрос ->
        # один основной queue.submit().
        # ---------------------------------------------------------------------

        main_messages: list[
            dict[str, Any]
        ] = [
            {
                "role": "system",
                "content": (
                    self.ai.get_system_prompt()
                ),
            }
        ]

        for item in history:
            if item.get("role") in {
                "user",
                "assistant",
            }:
                main_messages.append(
                    item
                )

        main_messages.append(
            {
                "role": "user",
                "content": main_user_content,
            }
        )

        answer = await self.queue.submit(
            main_messages
        )

        if not answer.strip():
            answer = (
                "Не удалось получить "
                "текстовый ответ."
            )

        # ---------------------------------------------------------------------
        # SAVE RAW HISTORY
        #
        # Сохраняем именно оригинальный ответ AI.
        # Не HTML-версию для Telegram.
        # ---------------------------------------------------------------------

        history.append(
            {
                "role": "user",
                "content": main_user_content,
            }
        )

        history.append(
            {
                "role": "assistant",
                "content": answer,
            }
        )

        # ---------------------------------------------------------------------
        # HISTORY COMPRESSION
        #
        # Это отдельный технический AI-вызов.
        # Он НЕ списывает пользовательский лимит.
        # ---------------------------------------------------------------------

        if (
            history_words(history)
            >= HISTORY_COMPRESS_WORDS
        ):
            history = (
                await self.compress_history(
                    history
                )
            )

        if len(history) > MAX_HISTORY_MESSAGES:
            history = history[
                -MAX_HISTORY_MESSAGES:
            ]

        await self.store.save_history(
            user_id,
            history,
        )

        # ---------------------------------------------------------------------
        # DISPLAY
        #
        # Нормализуем только то, что отправляется пользователю.
        # ---------------------------------------------------------------------

        display_answer = (
            normalize_ai_answer(
                answer
            )
        )

        if not display_answer:
            display_answer = (
                "Не удалось получить "
                "текстовый ответ."
            )

        try:
            await status_message.delete()
        except Exception:
            pass

        chunks = split_long_text(
            display_answer
        )

        for chunk in chunks:
            await message.answer(
                chunk
            )

    # =========================================================================
    # HISTORY COMPRESSION
    # =========================================================================

    async def compress_history(
        self,
        history: list[dict[str, Any]],
    ) -> list[dict[str, Any]]:

        if not history:
            return history

        compression_messages: list[
            dict[str, Any]
        ] = [
            {
                "role": "system",
                "content": (
                    self.ai.get_system_prompt()
                    + "\n\n"
                    "Ты выполняешь скрытое сжатие "
                    "истории диалога. "
                    "Пользователь не увидит этот запрос. "
                    "Сделай максимально полезное краткое "
                    "резюме предыдущего диалога: "
                    "сохрани факты, решения, контекст, "
                    "предпочтения пользователя, "
                    "незавершённые задачи и важные детали. "
                    "Не добавляй выдуманные сведения."
                ),
            }
        ]

        for item in history:
            role = item.get(
                "role"
            )

            if role not in {
                "user",
                "assistant",
            }:
                continue

            content = item.get(
                "content",
                "",
            )

            if isinstance(
                content,
                list,
            ):
                text_parts: list[str] = []

                for part in content:
                    if not isinstance(
                        part,
                        dict,
                    ):
                        continue

                    if part.get("type") == "text":
                        text = part.get(
                            "text",
                            "",
                        )

                        if isinstance(
                            text,
                            str,
                        ):
                            text_parts.append(
                                text
                            )

                content = "\n".join(
                    text_parts
                )

            if not isinstance(
                content,
                str,
            ):
                content = str(
                    content
                )

            compression_messages.append(
                {
                    "role": role,
                    "content": content,
                }
            )

        summary = await self.queue.submit(
            compression_messages
        )

        if not summary.strip():
            return history

        return [
            {
                "role": "assistant",
                "content": (
                    "[КРАТКОЕ РЕЗЮМЕ "
                    "ПРЕДЫДУЩЕГО ДИАЛОГА]\n"
                    + summary.strip()
                ),
            }
        ]


# =============================================================================
# GLOBAL APP
# =============================================================================

App = BotApp


# =============================================================================
# POLLING MODE
# =============================================================================

async def main() -> None:
    await App.start()

    try:
        await App.dp.start_polling(
            App.bot
        )
    finally:
        await App.stop()


# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        pass
